from __future__ import annotations

from dataclasses import dataclass
from typing import Any, Dict, List, Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE


@dataclass
class ParsedSlide:
    layout_group: str
    layout: str
    content: Dict[str, Any]
    speaker_note: Optional[str]
    properties: Dict[str, Any]


class PptxImportService:
    """Parse a .pptx file into the app's canonical slide representation."""

    @staticmethod
    def parse_pptx(path: str) -> List[ParsedSlide]:
        presentation = Presentation(path)
        parsed_slides: List[ParsedSlide] = []

        for index, slide in enumerate(presentation.slides):
            text_runs: List[Dict[str, Any]] = []
            images: List[Dict[str, Any]] = []

            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    for paragraph_index, paragraph in enumerate(shape.text_frame.paragraphs):
                        for run_index, run in enumerate(paragraph.runs):
                            text = (run.text or "").strip()
                            if not text:
                                continue

                            text_runs.append(
                                {
                                    "shape_id": getattr(shape, "shape_id", None),
                                    "shape_name": getattr(shape, "name", None),
                                    "paragraph_index": paragraph_index,
                                    "run_index": run_index,
                                    "text": text,
                                    "bold": run.font.bold,
                                    "italic": run.font.italic,
                                    "underline": run.font.underline,
                                    "font_name": run.font.name,
                                    "font_size": int(run.font.size.pt) if run.font.size else None,
                                    "font_color": PptxImportService._serialize_color(run.font.color.rgb),
                                }
                            )

                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image = shape.image
                    rel_id = getattr(shape, "_pic", None)
                    relationship_id = None
                    if rel_id is not None:
                        blip_fill = shape._pic.blipFill  # pylint: disable=protected-access
                        if blip_fill is not None and blip_fill.blip is not None:
                            relationship_id = blip_fill.blip.embed

                    images.append(
                        {
                            "shape_id": getattr(shape, "shape_id", None),
                            "shape_name": getattr(shape, "name", None),
                            "relationship_id": relationship_id,
                            "filename": image.filename,
                            "content_type": image.content_type,
                            "sha1": image.sha1,
                            "width": shape.width,
                            "height": shape.height,
                            "left": shape.left,
                            "top": shape.top,
                        }
                    )

            speaker_note = PptxImportService._extract_speaker_note(slide)
            slide_text = "\n".join(run["text"] for run in text_runs)
            slide_title = PptxImportService._extract_title(slide)

            parsed_slides.append(
                ParsedSlide(
                    layout_group="pptx-import",
                    layout=(slide.slide_layout.name or "unknown-layout").lower().replace(" ", "-"),
                    speaker_note=speaker_note,
                    properties={
                        "source": "pptx_import",
                        "slide_number": index + 1,
                        "layout_name": slide.slide_layout.name,
                    },
                    content={
                        "title": slide_title,
                        "text": slide_text,
                        "text_runs": text_runs,
                        "images": images,
                        "metadata": {
                            "source": "pptx_import",
                            "layout_name": slide.slide_layout.name,
                            "slide_number": index + 1,
                        },
                        "__speaker_note__": speaker_note,
                    },
                )
            )

        return parsed_slides

    @staticmethod
    def _extract_speaker_note(slide) -> Optional[str]:
        if not slide.has_notes_slide:
            return None

        notes_text = (slide.notes_slide.notes_text_frame.text or "").strip()
        return notes_text or None

    @staticmethod
    def _extract_title(slide) -> Optional[str]:
        if not slide.shapes.title:
            return None
        title_text = (slide.shapes.title.text or "").strip()
        return title_text or None

    @staticmethod
    def _serialize_color(color: Optional[RGBColor]) -> Optional[str]:
        if color is None:
            return None
        return str(color)
