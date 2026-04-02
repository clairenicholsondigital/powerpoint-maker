import zipfile
from pathlib import Path

import pytest
import xml.etree.ElementTree as ET
import base64
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


@pytest.fixture
def sample_assets(tmp_path: Path) -> dict[str, Path]:
    image_one = tmp_path / "image-one.png"
    image_two = tmp_path / "image-two.png"

    image_one.write_bytes(_png_bytes("iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAwMCAO7ZxwoAAAAASUVORK5CYII="))
    image_two.write_bytes(_png_bytes("iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8z8AABQEBhX5nVQAAAABJRU5ErkJggg=="))

    return {"image_one": image_one, "image_two": image_two}


@pytest.fixture
def sample_pptx(sample_assets: dict[str, Path], tmp_path: Path) -> Path:
    deck = Presentation()

    title_slide = deck.slides.add_slide(deck.slide_layouts[0])
    title_slide.shapes.title.text = "Quarterly Business Review"
    title_slide.placeholders[1].text = "Prepared for leadership"

    body_slide = deck.slides.add_slide(deck.slide_layouts[1])
    body_slide.shapes.title.text = "Growth Highlights"

    body = body_slide.shapes.placeholders[1].text_frame
    body.clear()

    first_para = body.paragraphs[0]
    run_a = first_para.add_run()
    run_a.text = "Revenue up 28% "
    run_a.font.name = "Calibri"
    run_a.font.size = Pt(22)
    run_a.font.color.rgb = RGBColor(0x1F, 0x4E, 0x78)

    run_b = first_para.add_run()
    run_b.text = "year-over-year"
    run_b.font.name = "Arial"
    run_b.font.size = Pt(20)
    run_b.font.bold = True
    run_b.font.color.rgb = RGBColor(0xC0, 0x00, 0x00)

    second_para = body.add_paragraph()
    second_para.text = "Pipeline quality improved across all regions"
    second_para.level = 1
    second_para.alignment = PP_ALIGN.LEFT

    body_slide.shapes.add_picture(str(sample_assets["image_one"]), Inches(6.2), Inches(1.6), width=Inches(2.6))
    body_slide.notes_slide.notes_text_frame.text = "Speaker note: Emphasize enterprise expansion."

    output = tmp_path / "sample-roundtrip.pptx"
    deck.save(output)
    return output


def _zip_parts(path: Path) -> dict[str, bytes]:
    with zipfile.ZipFile(path, "r") as archive:
        return {name: archive.read(name) for name in archive.namelist()}


def _canonical_xml_bytes(xml_bytes: bytes) -> bytes:
    root = ET.fromstring(xml_bytes)
    _sort_attrs_recursive(root)
    return ET.tostring(root, encoding="utf-8")


def _sort_attrs_recursive(element: ET.Element) -> None:
    if element.attrib:
        items = sorted(element.attrib.items())
        element.attrib.clear()
        element.attrib.update(items)
    for child in list(element):
        _sort_attrs_recursive(child)


def test_import_fixture_pptx(sample_pptx: Path):
    presentation = Presentation(str(sample_pptx))

    assert len(presentation.slides) == 2
    assert presentation.slides[0].shapes.title.text == "Quarterly Business Review"

    slide = presentation.slides[1]
    assert slide.shapes.title.text == "Growth Highlights"
    assert slide.notes_slide.notes_text_frame.text == "Speaker note: Emphasize enterprise expansion."

    picture_shapes = [shape for shape in slide.shapes if shape.shape_type == 13]
    assert picture_shapes, "Expected at least one image in fixture PPTX"


def test_edit_and_round_trip_export_with_xml_deltas(sample_pptx: Path, sample_assets: dict[str, Path], tmp_path: Path):
    before_parts = _zip_parts(sample_pptx)

    presentation = Presentation(str(sample_pptx))
    target_slide = presentation.slides[1]

    target_slide.shapes.title.text = "Growth Highlights (Revised)"
    target_slide.shapes.placeholders[1].text_frame.paragraphs[0].runs[0].text = "Revenue up 31% "
    target_slide.notes_slide.notes_text_frame.text = "Speaker note: Highlight updated guidance."

    old_picture = next(shape for shape in target_slide.shapes if shape.shape_type == 13)
    left, top, width, height = old_picture.left, old_picture.top, old_picture.width, old_picture.height
    old_picture._element.getparent().remove(old_picture._element)
    target_slide.shapes.add_picture(str(sample_assets["image_two"]), left, top, width=width, height=height)

    exported = tmp_path / "roundtrip-exported.pptx"
    presentation.save(exported)

    reopened = Presentation(str(exported))
    assert reopened.slides[1].shapes.title.text == "Growth Highlights (Revised)"
    assert reopened.slides[1].notes_slide.notes_text_frame.text == "Speaker note: Highlight updated guidance."

    after_parts = _zip_parts(exported)

    assert b"Growth Highlights (Revised)" in after_parts["ppt/slides/slide2.xml"]
    assert b"Revenue up 31%" in after_parts["ppt/slides/slide2.xml"]
    notes_parts = [name for name in after_parts if name.startswith("ppt/notesSlides/") and name.endswith(".xml") and "/_rels/" not in name]
    assert notes_parts
    assert any(b"Highlight updated guidance" in after_parts[name] for name in notes_parts)

    # Untouched slide XML stays semantically equivalent across round trip.
    assert _canonical_xml_bytes(before_parts["ppt/slides/slide1.xml"]) == _canonical_xml_bytes(
        after_parts["ppt/slides/slide1.xml"]
    )

    # For media binaries, byte-stability is not feasible after image replacement.
    # Verify untouched image assets remain valid PNGs where present.
    media_parts = [name for name in after_parts if name.startswith("ppt/media/")]
    assert media_parts
    for part_name in media_parts:
        assert after_parts[part_name].startswith(PNG_SIGNATURE)


PNG_SIGNATURE = b"\x89PNG\r\n\x1a\n"


def _png_bytes(data: str) -> bytes:
    return base64.b64decode(data)
